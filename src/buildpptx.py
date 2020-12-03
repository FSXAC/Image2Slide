import os

from pptx import Presentation
from pptx.util import Inches

USE_IMG_BG = False

p = Presentation()
p.slide_height = Inches(9)
p.slide_width = Inches(16)

TITLE_LAYOUT = p.slide_layouts[0]
CONTENT_LAYOUT = p.slide_layouts[8]

# slide = p.slides.add_slide(TITLE_LAYOUT)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]

# title.text = 'Hello world'
# subtitle.text = 'EECE 570 Project Report'

# ==

# slide = p.slides.add_slide(CONTENT_LAYOUT)
# title = slide.shapes.title
# title.text = 'Image'
# pic = slide.shapes.add_picture('./src/img/ipad.jpg', Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

# ==

# Takes all the images in output folder and sticks them together
# for (dirpath, dirnames, filenames) in os.walk('src/output'):
#     print(dirpath, dirnames, filenames)

root_path = './output'
for fitem in sorted(os.listdir(root_path)):
    sub_path = os.path.join(root_path, fitem)

    if os.path.isdir(sub_path):

        slide = p.slides.add_slide(CONTENT_LAYOUT)

        # add background
        if (USE_IMG_BG):
            slide.shapes.add_picture(
                os.path.join(sub_path, 'bg.jpg'),
                0, 0, p.slide_width, p.slide_height
            )

        with open(os.path.join(sub_path, 'meta.csv')) as metacsv:
            # meta.csv contains all the regions in x,y,w,h form
            roi_metadata = metacsv.readlines()

            for i, roi in enumerate(roi_metadata):
                data = [float(x) for x in roi.split(',')]
                x, y, w, h = data

                slide.shapes.add_picture(
                    os.path.join(sub_path, f'{i + 1}.png'),
                    x * p.slide_height,
                    y * p.slide_height,
                    w * p.slide_height,
                    h * p.slide_height
                )
            
p.save('../test.pptx')
